from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "What is an NEA?"
subtitle1.text = "My A-Level Computer Science Project\n\nPresenter: Mohid\nAQA A-Level Computer Science"

notes1 = slide1.notes_slide.notes_text_frame
notes1.text = """Hello everyone! Today I'm going to talk about my A-Level Computer Science project. I know most of you probably haven't done Computer Science yet, so I'll start by explaining what an NEA actually is. NEA stands for Non-Exam Assessment - basically, it's a big coursework project that makes up 20% of your final A-Level grade. Instead of just sitting exams, you get to spend months building a real piece of software that solves an actual problem. Think of it like your coursework in other subjects, but instead of writing essays, you're creating working computer programs."""

# Slide 2: What is an NEA?
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
title2.text = "What is an NEA?"

content2 = slide2.placeholders[1].text_frame
content2.text = "The Non-Exam Assessment Explained"
p = content2.add_paragraph()
p.text = "20% of your final A-Level grade"
p.level = 1
p = content2.add_paragraph()
p.text = "A real programming project - not just theory!"
p.level = 1
p = content2.add_paragraph()
p.text = "4 main stages:"
p.level = 1
p = content2.add_paragraph()
p.text = "Analysis: What problem are you solving?"
p.level = 2
p = content2.add_paragraph()
p.text = "Design: How will you build it?"
p.level = 2
p = content2.add_paragraph()
p.text = "Development: Actually coding it"
p.level = 2
p = content2.add_paragraph()
p.text = "Evaluation: Did it work?"
p.level = 2

notes2 = slide2.notes_slide.notes_text_frame
notes2.text = """So what exactly is the NEA? In AQA Computer Science, the NEA is worth 20% of your final grade, which is pretty significant. Unlike traditional exams where you answer questions, the NEA requires you to actually build a real piece of software from scratch. You go through four main stages: First, Analysis - where you identify a problem that needs solving and figure out what your program needs to do. Second, Design - where you plan out exactly how you'll build it, what features it needs, and how users will interact with it. Third, Development - this is the actual coding part where you spend months writing the program. And finally, Evaluation - where you test everything, fix bugs, and see if it actually solved the problem you set out to fix. The whole process takes around 6-9 months from start to finish."""

# Slide 3: My Project
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
title3.text = "My Project - Scientific Calculator"

content3 = slide3.placeholders[1].text_frame
content3.text = "Building a Calculator for Engineers & Scientists"
p = content3.add_paragraph()
p.text = "Problem: Students need a scientific calculator for A-Level Maths & Physics"
p.level = 1
p = content3.add_paragraph()
p.text = "Solution: A custom calculator that:"
p.level = 1
p = content3.add_paragraph()
p.text = "Works on physical hardware (not just a computer screen)"
p.level = 2
p = content3.add_paragraph()
p.text = "Has advanced features like graphing and equation solving"
p.level = 2
p = content3.add_paragraph()
p.text = "Looks and feels like a real Casio calculator"
p.level = 2

notes3 = slide3.notes_slide.notes_text_frame
notes3.text = """For my NEA, I decided to build a scientific calculator. You might be thinking "don't calculators already exist?" - and you'd be right! But here's the thing: I wanted to create something that works on actual physical hardware that you can hold in your hand, not just software on a computer. The problem I identified was that students doing A-Level Maths and Physics need calculators with advanced features - things like graphing functions, solving equations, and doing statistical calculations. My goal was to create a calculator that looks and feels like the Casio calculators you use in class, but is completely custom-built from scratch. This means I had to program everything - from how buttons work, to how numbers are displayed on screen, to the actual mathematics behind scientific calculations."""

# Slide 4: Technology
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
title4.text = "The Technology Behind It"

content4 = slide4.placeholders[1].text_frame
content4.text = "How Does It Actually Work?"
p = content4.add_paragraph()
p.text = "Hardware: ESP32 microcontroller (a tiny computer chip)"
p.level = 1
p = content4.add_paragraph()
p.text = "Display: LCD screen (like your phone, but simpler)"
p.level = 1
p = content4.add_paragraph()
p.text = "Buttons: Physical keypad (like a real calculator)"
p.level = 1
p = content4.add_paragraph()
p.text = "Programming Language: C programming"
p.level = 1
p = content4.add_paragraph()
p.text = "Lines of Code: Over 8,000 lines!"
p.level = 1

notes4 = slide4.notes_slide.notes_text_frame
notes4.text = """Let me explain the technology that makes this work - don't worry, I'll keep it simple! The "brain" of the calculator is called an ESP32 - it's a tiny computer chip about the size of a coin that you can program to do whatever you want. It has a small LCD screen to display numbers and graphs, and a physical button keypad that you press just like a normal calculator. I programmed it using a language called C, which is one of the fundamental programming languages used in real-world devices like cars, phones, and appliances. By the time I finished, I'd written over 8,000 lines of code - that's like writing a 100-page essay, except every single word has to be perfect or the whole thing breaks! Each line tells the computer exactly what to do, whether that's "when button 5 is pressed, display a 5" or "calculate the sine of this angle."."""

# Slide 5: Features
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
title5.text = "Cool Features I Built"

content5 = slide5.placeholders[1].text_frame
content5.text = "What Can It Actually Do?"
p = content5.add_paragraph()
p.text = "Math Mode: Basic and advanced calculations (sin, cos, powers, roots)"
p.level = 1
p = content5.add_paragraph()
p.text = "Graph Mode: Plot multiple functions on screen"
p.level = 1
p = content5.add_paragraph()
p.text = "Solver Mode: Solve equations automatically"
p.level = 1
p = content5.add_paragraph()
p.text = "Statistics Mode: Calculate probabilities and distributions"
p.level = 1
p = content5.add_paragraph()
p.text = "Mechanics Mode: Solve physics motion problems (SUVAT)"
p.level = 1

notes5 = slide5.notes_slide.notes_text_frame
notes5.text = """So what can this calculator actually do? Well, it has five main apps or modes. Math Mode is your basic calculator - you can do addition, multiplication, but also advanced stuff like trigonometry, powers, and square roots. Graph Mode lets you type in mathematical functions and it draws the graphs on screen - you can zoom in, zoom out, and even trace along the curve to see exact values. Solver Mode is really cool - you can type in an equation and it will solve it for you automatically using something called Newton-Raphson method, which is an algorithm for finding solutions. Statistics Mode helps with probability - you can calculate things like normal distributions and binomial probabilities, which are crucial for A-Level Maths. And finally, Mechanics Mode solves physics SUVAT equations for you - that's the motion equations you learn in A-Level Physics with displacement, velocity, acceleration, and time. The calculator figures out the unknown values based on what you give it."""

# Slide 6: Challenges
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title6 = slide6.shapes.title
title6.text = "Challenges & Problem Solving"

content6 = slide6.placeholders[1].text_frame
content6.text = "The Hard Parts"
p = content6.add_paragraph()
p.text = "Challenges I faced:"
p.level = 1
p = content6.add_paragraph()
p.text = "Making buttons work reliably (debouncing)"
p.level = 2
p = content6.add_paragraph()
p.text = "Creating custom fonts for mathematical symbols (π, √, ²)"
p.level = 2
p = content6.add_paragraph()
p.text = "Building a reusable menu system across different apps"
p.level = 2
p = content6.add_paragraph()
p.text = "Fixing the graph editor keyboard input bugs"
p.level = 2
p = content6.add_paragraph()
p.text = "Solution approach:"
p.level = 1
p = content6.add_paragraph()
p.text = "Test, debug, fix - repeat hundreds of times!"
p.level = 2
p = content6.add_paragraph()
p.text = "Research how real calculators handle these problems"
p.level = 2
p = content6.add_paragraph()
p.text = "Break big problems into smaller pieces"
p.level = 2

notes6 = slide6.notes_slide.notes_text_frame
notes6.text = """Building this wasn't easy - I ran into loads of challenges along the way. One of the first problems was making the buttons work reliably. When you press a physical button, it actually "bounces" and sends multiple signals - I had to write code to filter that out, which is called debouncing. Another big challenge was creating custom fonts to display mathematical symbols like pi, square roots, and superscript numbers - computers don't just "know" how to show these, you have to design every single pixel. Later in the project, I had to build a reusable menu system so that all the different calculator modes had consistent navigation - this meant planning out a design pattern that could work across different apps. One of the trickiest bugs to fix was in the graph editor - when you pressed the 'x' key, it would sometimes type 'V' instead! This took days to debug, but I eventually figured out it was because keyboard events were being processed twice. The key to solving all these problems was breaking them down into smaller pieces, testing constantly, and not giving up when things went wrong."""

# Slide 7: Conclusion
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title7 = slide7.shapes.title
title7.text = "Key Takeaways & Conclusion"

content7 = slide7.placeholders[1].text_frame
content7.text = "What I Learned"
p = content7.add_paragraph()
p.text = "Skills developed:"
p.level = 1
p = content7.add_paragraph()
p.text = "Programming in C (low-level hardware control)"
p.level = 2
p = content7.add_paragraph()
p.text = "Problem solving & debugging"
p.level = 2
p = content7.add_paragraph()
p.text = "Project management over 6+ months"
p.level = 2
p = content7.add_paragraph()
p.text = "User interface design"
p.level = 2
p = content7.add_paragraph()
p.text = "Why this matters:"
p.level = 1
p = content7.add_paragraph()
p.text = "Real-world software development experience"
p.level = 2
p = content7.add_paragraph()
p.text = "Understanding how everyday devices actually work"
p.level = 2
p = content7.add_paragraph()
p.text = "Preparing for university Computer Science courses"
p.level = 2
p = content7.add_paragraph()
p.text = ""
p.level = 1
p = content7.add_paragraph()
p.text = "Questions?"
p.level = 0

notes7 = slide7.notes_slide.notes_text_frame
notes7.text = """So what did I actually learn from doing this NEA? Well, obviously I got much better at programming in C, which is a really powerful language used in embedded systems - that's devices like calculators, cars, medical equipment, and smart appliances. But beyond just coding, I learned how to debug complex problems, which means figuring out why something isn't working and fixing it. I also learned project management - keeping track of what needs to be done, prioritizing features, and managing my time over six months. And I learned about user interface design - making sure the calculator is actually pleasant and intuitive to use, not just functional. Why does this matter? Well, this NEA gives you real-world software development experience that universities and employers actually value. You're not just learning theory from textbooks - you're building something real that actually works. It helps you understand how the devices you use every day actually function under the hood. And if you're thinking about studying Computer Science at university, the NEA is excellent preparation because it teaches you the problem-solving approach that professional developers use every single day. Thanks for listening - does anyone have any questions?"""

# Save presentation
prs.save('NEA_Presentation.pptx')
print("Presentation created successfully: NEA_Presentation.pptx")
