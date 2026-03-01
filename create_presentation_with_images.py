from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
import io

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Helper function to create simple diagram images
def create_nea_stages_image():
    """Create a diagram showing the 4 NEA stages"""
    img = Image.new('RGB', (800, 400), color='white')
    draw = ImageDraw.Draw(img)
    
    # Draw 4 boxes for stages
    stages = ['Analysis', 'Design', 'Development', 'Evaluation']
    colors = ['#3498db', '#e74c3c', '#2ecc71', '#f39c12']
    
    box_width = 150
    box_height = 80
    spacing = 40
    start_x = 50
    y = 160
    
    for i, (stage, color) in enumerate(zip(stages, colors)):
        x = start_x + i * (box_width + spacing)
        # Draw box
        draw.rectangle([x, y, x + box_width, y + box_height], fill=color, outline='black', width=3)
        # Draw text (simplified - no font for simplicity)
        text_x = x + box_width // 2
        text_y = y + box_height // 2
        draw.text((text_x - 30, text_y - 10), stage, fill='white')
        
        # Draw arrow
        if i < 3:
            arrow_x = x + box_width + 10
            draw.polygon([(arrow_x, y + 40), (arrow_x + 20, y + 40), 
                         (arrow_x + 10, y + 50), (arrow_x + 20, y + 40),
                         (arrow_x + 20, y + 30), (arrow_x + 10, y + 40)], 
                        fill='#555555')
    
    # Save to bytes
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    return img_bytes

def create_calculator_diagram():
    """Create a simple calculator hardware diagram"""
    img = Image.new('RGB', (600, 500), color='white')
    draw = ImageDraw.Draw(img)
    
    # Draw calculator outline
    draw.rectangle([100, 50, 500, 450], fill='#2c3e50', outline='black', width=3)
    
    # Draw screen
    draw.rectangle([130, 80, 470, 180], fill='#ecf0f1', outline='black', width=2)
    draw.text((180, 120), "3.14159...", fill='black')
    
    # Draw button grid
    button_colors = ['#34495e', '#e74c3c', '#3498db', '#2ecc71']
    for row in range(4):
        for col in range(4):
            x = 140 + col * 80
            y = 220 + row * 50
            color = button_colors[col % 4]
            draw.rectangle([x, y, x + 60, y + 35], fill=color, outline='black', width=2)
    
    # Add labels
    draw.text((250, 20), "ESP32 Calculator", fill='black')
    
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    return img_bytes

def create_features_icons():
    """Create simple icons for the 5 calculator modes"""
    img = Image.new('RGB', (800, 200), color='white')
    draw = ImageDraw.Draw(img)
    
    features = ['Math', 'Graph', 'Solver', 'Stats', 'Mech']
    colors = ['#3498db', '#e74c3c', '#f39c12', '#2ecc71', '#9b59b6']
    
    icon_size = 120
    spacing = 20
    start_x = 40
    y = 40
    
    for i, (feature, color) in enumerate(zip(features, colors)):
        x = start_x + i * (icon_size + spacing)
        # Draw circle icon
        draw.ellipse([x, y, x + icon_size, y + icon_size], fill=color, outline='black', width=2)
        # Draw text
        draw.text((x + 30, y + 50), feature, fill='white')
    
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    return img_bytes

def create_challenge_diagram():
    """Create a problem-solving cycle diagram"""
    img = Image.new('RGB', (500, 500), color='white')
    draw = ImageDraw.Draw(img)
    
    # Draw circular problem-solving cycle
    center_x, center_y = 250, 250
    radius = 150
    
    steps = ['Problem', 'Debug', 'Test', 'Fix']
    colors = ['#e74c3c', '#f39c12', '#3498db', '#2ecc71']
    
    for i, (step, color) in enumerate(zip(steps, colors)):
        angle = i * 90
        x = center_x + int(radius * 0.7 * (1 if i % 2 == 0 else -1) * (1 if i < 2 else 1))
        y = center_y + int(radius * 0.7 * (1 if i < 2 else -1) * (1 if i % 2 == 1 else 1))
        
        # Draw box
        draw.rectangle([x - 60, y - 30, x + 60, y + 30], fill=color, outline='black', width=3)
        draw.text((x - 25, y - 10), step, fill='white')
    
    # Draw center circle
    draw.ellipse([center_x - 40, center_y - 40, center_x + 40, center_y + 40], 
                 fill='#34495e', outline='black', width=3)
    draw.text((center_x - 25, center_y - 10), "Repeat", fill='white')
    
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    return img_bytes

# Slide 1: Title (no image needed)
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "What is an NEA?"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "My A-Level Computer Science Project\n\nPresenter: Mohid\nAQA A-Level Computer Science"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(24)
subtitle_para.alignment = PP_ALIGN.CENTER

notes1 = slide1.notes_slide.notes_text_frame
notes1.text = """Hello everyone! Today I'm going to talk about my A-Level Computer Science project. I know most of you probably haven't done Computer Science yet, so I'll start by explaining what an NEA actually is. NEA stands for Non-Exam Assessment - basically, it's a big coursework project that makes up 20% of your final A-Level grade. Instead of just sitting exams, you get to spend months building a real piece of software that solves an actual problem. Think of it like your coursework in other subjects, but instead of writing essays, you're creating working computer programs."""

# Slide 2: What is an NEA? (with stages diagram)
slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
title2 = slide2.shapes.title
title2.text = "What is an NEA?"

# Add text box for bullet points
text_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
text_frame = text_box.text_frame
text_frame.text = "The Non-Exam Assessment"
p = text_frame.add_paragraph()
p.text = "20% of your final grade"
p.level = 1
p = text_frame.add_paragraph()
p.text = "A real programming project"
p.level = 1
p = text_frame.add_paragraph()
p.text = "4 main stages:"
p.level = 1
p = text_frame.add_paragraph()
p.text = "Analysis"
p.level = 2
p = text_frame.add_paragraph()
p.text = "Design"
p.level = 2
p = text_frame.add_paragraph()
p.text = "Development"
p.level = 2
p = text_frame.add_paragraph()
p.text = "Evaluation"
p.level = 2

# Add diagram image
img_stream = create_nea_stages_image()
pic = slide2.shapes.add_picture(img_stream, Inches(5.2), Inches(2), width=Inches(4.5))

notes2 = slide2.notes_slide.notes_text_frame
notes2.text = """So what exactly is the NEA? In AQA Computer Science, the NEA is worth 20% of your final grade, which is pretty significant. Unlike traditional exams where you answer questions, the NEA requires you to actually build a real piece of software from scratch. You go through four main stages: First, Analysis - where you identify a problem that needs solving and figure out what your program needs to do. Second, Design - where you plan out exactly how you'll build it, what features it needs, and how users will interact with it. Third, Development - this is the actual coding part where you spend months writing the program. And finally, Evaluation - where you test everything, fix bugs, and see if it actually solved the problem you set out to fix. The whole process takes around 6-9 months from start to finish."""

# Slide 3: My Project (with calculator image)
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
title3 = slide3.shapes.title
title3.text = "My Project - Scientific Calculator"

text_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
text_frame3 = text_box3.text_frame
text_frame3.text = "Problem: Students need scientific calculators"
p = text_frame3.add_paragraph()
p.text = "Solution: Custom calculator"
p.level = 1
p = text_frame3.add_paragraph()
p.text = "Physical hardware"
p.level = 2
p = text_frame3.add_paragraph()
p.text = "Advanced features"
p.level = 2
p = text_frame3.add_paragraph()
p.text = "Graphing"
p.level = 2
p = text_frame3.add_paragraph()
p.text = "Equation solving"
p.level = 2
p = text_frame3.add_paragraph()
p.text = "Statistics"
p.level = 2

# Add calculator diagram
img_stream3 = create_calculator_diagram()
pic3 = slide3.shapes.add_picture(img_stream3, Inches(5.5), Inches(1.8), width=Inches(4))

notes3 = slide3.notes_slide.notes_text_frame
notes3.text = """For my NEA, I decided to build a scientific calculator. You might be thinking "don't calculators already exist?" - and you'd be right! But here's the thing: I wanted to create something that works on actual physical hardware that you can hold in your hand, not just software on a computer. The problem I identified was that students doing A-Level Maths and Physics need calculators with advanced features - things like graphing functions, solving equations, and doing statistical calculations. My goal was to create a calculator that looks and feels like the Casio calculators you use in class, but is completely custom-built from scratch. This means I had to program everything - from how buttons work, to how numbers are displayed on screen, to the actual mathematics behind scientific calculations."""

# Slide 4: Technology (with hardware components)
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
title4.text = "The Technology Behind It"

content4 = slide4.placeholders[1].text_frame
content4.text = "How Does It Work?"
p = content4.add_paragraph()
p.text = "ESP32 microcontroller (tiny computer chip)"
p.level = 1
p = content4.add_paragraph()
p.text = "LCD screen"
p.level = 1
p = content4.add_paragraph()
p.text = "Physical keypad"
p.level = 1
p = content4.add_paragraph()
p.text = "C programming language"
p.level = 1
p = content4.add_paragraph()
p.text = "Over 8,000 lines of code!"
p.level = 1

notes4 = slide4.notes_slide.notes_text_frame
notes4.text = """Let me explain the technology that makes this work - don't worry, I'll keep it simple! The "brain" of the calculator is called an ESP32 - it's a tiny computer chip about the size of a coin that you can program to do whatever you want. It has a small LCD screen to display numbers and graphs, and a physical button keypad that you press just like a normal calculator. I programmed it using a language called C, which is one of the fundamental programming languages used in real-world devices like cars, phones, and appliances. By the time I finished, I'd written over 8,000 lines of code - that's like writing a 100-page essay, except every single word has to be perfect or the whole thing breaks! Each line tells the computer exactly what to do, whether that's "when button 5 is pressed, display a 5" or "calculate the sine of this angle."."""

# Slide 5: Features (with icons)
slide5 = prs.slides.add_slide(prs.slide_layouts[5])
title5 = slide5.shapes.title
title5.text = "Cool Features I Built"

text_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
text_frame5 = text_box5.text_frame
text_frame5.text = "What Can It Actually Do?"

# Add feature icons
img_stream5 = create_features_icons()
pic5 = slide5.shapes.add_picture(img_stream5, Inches(1), Inches(3), width=Inches(8))

# Add descriptions below
desc_box = slide5.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
desc_frame = desc_box.text_frame
desc_frame.text = "Math • Graphing • Equation Solving • Statistics • Physics (SUVAT)"
desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

notes5 = slide5.notes_slide.notes_text_frame
notes5.text = """So what can this calculator actually do? Well, it has five main apps or modes. Math Mode is your basic calculator - you can do addition, multiplication, but also advanced stuff like trigonometry, powers, and square roots. Graph Mode lets you type in mathematical functions and it draws the graphs on screen - you can zoom in, zoom out, and even trace along the curve to see exact values. Solver Mode is really cool - you can type in an equation and it will solve it for you automatically using something called Newton-Raphson method, which is an algorithm for finding solutions. Statistics Mode helps with probability - you can calculate things like normal distributions and binomial probabilities, which are crucial for A-Level Maths. And finally, Mechanics Mode solves physics SUVAT equations for you - that's the motion equations you learn in A-Level Physics with displacement, velocity, acceleration, and time. The calculator figures out the unknown values based on what you give it."""

# Slide 6: Challenges (with problem-solving cycle)
slide6 = prs.slides.add_slide(prs.slide_layouts[5])
title6 = slide6.shapes.title
title6.text = "Challenges & Problem Solving"

text_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
text_frame6 = text_box6.text_frame
text_frame6.text = "Challenges:"
p = text_frame6.add_paragraph()
p.text = "Button debouncing"
p.level = 1
p = text_frame6.add_paragraph()
p.text = "Custom fonts (π, √, ²)"
p.level = 1
p = text_frame6.add_paragraph()
p.text = "Menu system design"
p.level = 1
p = text_frame6.add_paragraph()
p.text = "Keyboard input bugs"
p.level = 1
p = text_frame6.add_paragraph()
p.text = "Solution: Test, debug, repeat!"
p.level = 0

# Add problem-solving diagram
img_stream6 = create_challenge_diagram()
pic6 = slide6.shapes.add_picture(img_stream6, Inches(5.5), Inches(2), width=Inches(4))

notes6 = slide6.notes_slide.notes_text_frame
notes6.text = """Building this wasn't easy - I ran into loads of challenges along the way. One of the first problems was making the buttons work reliably. When you press a physical button, it actually "bounces" and sends multiple signals - I had to write code to filter that out, which is called debouncing. Another big challenge was creating custom fonts to display mathematical symbols like pi, square roots, and superscript numbers - computers don't just "know" how to show these, you have to design every single pixel. Later in the project, I had to build a reusable menu system so that all the different calculator modes had consistent navigation - this meant planning out a design pattern that could work across different apps. One of the trickiest bugs to fix was in the graph editor - when you pressed the 'x' key, it would sometimes type 'V' instead! This took days to debug, but I eventually figured out it was because keyboard events were being processed twice. The key to solving all these problems was breaking them down into smaller pieces, testing constantly, and not giving up when things went wrong."""

# Slide 7: Conclusion
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title7 = slide7.shapes.title
title7.text = "Key Takeaways"

content7 = slide7.placeholders[1].text_frame
content7.text = "What I Learned:"
p = content7.add_paragraph()
p.text = "C programming & hardware control"
p.level = 1
p = content7.add_paragraph()
p.text = "Problem solving & debugging"
p.level = 1
p = content7.add_paragraph()
p.text = "Project management"
p.level = 1
p = content7.add_paragraph()
p.text = "UI/UX design"
p.level = 1
p = content7.add_paragraph()
p.text = ""
p = content7.add_paragraph()
p.text = "Why it matters:"
p.level = 0
p = content7.add_paragraph()
p.text = "Real-world experience"
p.level = 1
p = content7.add_paragraph()
p.text = "Prepare for university"
p.level = 1
p = content7.add_paragraph()
p.text = ""
p = content7.add_paragraph()
p.text = "Questions?"
p.level = 0
p.font.size = Pt(32)
p.font.bold = True

notes7 = slide7.notes_slide.notes_text_frame
notes7.text = """So what did I actually learn from doing this NEA? Well, obviously I got much better at programming in C, which is a really powerful language used in embedded systems - that's devices like calculators, cars, medical equipment, and smart appliances. But beyond just coding, I learned how to debug complex problems, which means figuring out why something isn't working and fixing it. I also learned project management - keeping track of what needs to be done, prioritizing features, and managing my time over six months. And I learned about user interface design - making sure the calculator is actually pleasant and intuitive to use, not just functional. Why does this matter? Well, this NEA gives you real-world software development experience that universities and employers actually value. You're not just learning theory from textbooks - you're building something real that actually works. It helps you understand how the devices you use every day actually function under the hood. And if you're thinking about studying Computer Science at university, the NEA is excellent preparation because it teaches you the problem-solving approach that professional developers use every single day. Thanks for listening - does anyone have any questions?"""

# Save presentation
prs.save('NEA_Presentation_With_Images.pptx')
print("Presentation with images created successfully: NEA_Presentation_With_Images.pptx")
